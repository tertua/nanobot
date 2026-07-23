"""Document text extraction utilities for nanobot."""

import mimetypes
from collections.abc import Iterator
from dataclasses import dataclass
from pathlib import Path
from zipfile import BadZipFile, ZipFile

from loguru import logger

from nanobot.utils.helpers import detect_image_mime

# Supported file extensions for text extraction
SUPPORTED_EXTENSIONS: set[str] = {
    # Document formats
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    # Text formats
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".log",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    # Image formats (for future OCR support)
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
}

_MAX_TEXT_LENGTH = 200_000
_MAX_EXTRACT_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
_MAX_OFFICE_ARCHIVE_MEMBERS = 10_000
_MAX_OFFICE_UNCOMPRESSED_SIZE = 256 * 1024 * 1024  # 256 MB
_MAX_OFFICE_MEMBER_SIZE = 128 * 1024 * 1024  # 128 MB
_MAX_DOCX_TABLE_CELLS = 100_000
_MAX_DOCX_TABLE_DEPTH = 8
_MAX_PDF_CONTENT_STREAM_SIZE = 32 * 1024 * 1024  # 32 MB per page
_MAX_PDF_ATTACHMENT_PAGES = 100


class _TextCollector:
    """Build bounded parser output without retaining the full document text."""

    def __init__(self, limit: int) -> None:
        self.limit = limit
        self.parts: list[str] = []
        self.length = 0
        self.truncated = False

    def add(self, text: str, *, separator: str = "") -> bool:
        if not text:
            return True
        prefix = separator if self.parts else ""
        chunk = prefix + text
        remaining = self.limit - self.length
        if len(chunk) > remaining:
            if remaining > 0:
                self.parts.append(chunk[:remaining])
                self.length += remaining
            self.truncated = True
            return False
        self.parts.append(chunk)
        self.length += len(chunk)
        return True

    def render(self) -> str:
        text = "".join(self.parts)
        if self.truncated:
            text += f"... (truncated at {self.limit} chars)"
        return text


class PdfSafetyError(Exception):
    """Raised when a PDF exceeds a parser safety boundary."""


class PdfPageRangeError(Exception):
    """Raised when a requested PDF page range is invalid."""


class DocxSafetyError(Exception):
    """Raised when a DOCX table exceeds a parser safety boundary."""


@dataclass(frozen=True, slots=True)
class PdfExtraction:
    text: str
    total_pages: int
    start_page: int
    end_page: int


def extract_text(path: Path) -> str | None:
    """Extract text from a file.

    Args:
        path: Path to the file.

    Returns:
        Extracted text as string, None for unsupported types,
        or error string for failures.
    """
    if not isinstance(path, Path):
        path = Path(path)

    if not path.exists():
        return f"[error: file not found: {path}]"
    try:
        if path.stat().st_size > _MAX_EXTRACT_FILE_SIZE:
            return f"[error: file exceeds {_MAX_EXTRACT_FILE_SIZE // (1024 * 1024)} MB limit]"
    except OSError as e:
        return f"[error: failed to inspect file: {e!s}]"

    ext = path.suffix.lower()

    # Parsers stay lazy even though they are bundled so idle processes do not
    # retain their import cost (see issue #3422).
    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext == ".docx":
        return _extract_docx(path)
    elif ext == ".xlsx":
        return _extract_xlsx(path)
    elif ext == ".pptx":
        return _extract_pptx(path)
    elif _is_text_extension(ext):
        return _extract_text_file(path)
    elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        # Image files - for future OCR support
        return f"[image: {path.name}]"
    else:
        # Unsupported extension
        return None


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using pypdf."""
    try:
        result = extract_pdf_pages(
            path,
            max_pages=_MAX_PDF_ATTACHMENT_PAGES,
            max_chars=_MAX_TEXT_LENGTH,
        )
        text = result.text
        if result.end_page < result.total_pages - 1:
            text += f"\n\n(Showing pages 1-{result.end_page + 1} of {result.total_pages}.)"
        return text
    except Exception as e:
        logger.exception("Failed to extract PDF {}", path)
        return f"[error: failed to extract PDF: {e!s}]"


def extract_pdf_pages(
    path: Path,
    *,
    pages: str | None = None,
    max_pages: int = _MAX_PDF_ATTACHMENT_PAGES,
    max_chars: int = _MAX_TEXT_LENGTH,
) -> PdfExtraction:
    """Extract a bounded PDF page range using the bundled pypdf reader."""
    from pypdf import PdfReader

    reader = PdfReader(path, strict=False)
    total_pages = len(reader.pages)
    if total_pages == 0:
        return PdfExtraction("", 0, 0, -1)

    start, end = _parse_pdf_page_range(pages, total_pages)
    end = min(end, start + max_pages - 1)
    collector = _TextCollector(max_chars)
    for index in range(start, end + 1):
        page = reader.pages[index]
        contents = page.get_contents()
        if contents is not None:
            stream_size = len(contents.get_data())
            if stream_size > _MAX_PDF_CONTENT_STREAM_SIZE:
                raise PdfSafetyError(
                    f"page {index + 1} content stream exceeds "
                    f"{_MAX_PDF_CONTENT_STREAM_SIZE // (1024 * 1024)} MB limit"
                )
        text = (page.extract_text() or "").strip()
        if text and not collector.add(f"--- Page {index + 1} ---\n{text}", separator="\n\n"):
            end = index
            break
    return PdfExtraction(collector.render(), total_pages, start, end)


def _parse_pdf_page_range(pages: str | None, total_pages: int) -> tuple[int, int]:
    if not pages:
        return 0, total_pages - 1
    values = pages.strip().split("-")
    if len(values) not in {1, 2}:
        raise PdfPageRangeError(f"invalid page range: {pages}")
    try:
        start = int(values[0])
        end = int(values[-1])
    except ValueError as e:
        raise PdfPageRangeError(f"invalid page range: {pages}") from e
    if start < 1 or end < start or start > total_pages:
        raise PdfPageRangeError(f"invalid page range: {pages}")
    return start - 1, min(end, total_pages) - 1


def _extract_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document as DocxDocument
        from docx.table import Table, _Cell
        from docx.text.paragraph import Paragraph
    except ImportError:
        return "[error: python-docx not installed]"
    try:
        if error := _office_archive_error(path):
            return error
        doc = DocxDocument(path)
        collector = _TextCollector(_MAX_TEXT_LENGTH)
        table_cell_count = 0

        def cell_text(cell: _Cell, depth: int) -> str:
            parts: list[str] = []
            for block in cell.iter_inner_content():
                if isinstance(block, Paragraph):
                    text = " ".join(block.text.split())
                    if text:
                        parts.append(text)
                elif isinstance(block, Table):
                    parts.extend(row.replace("\t", " | ") for row in table_rows(block, depth + 1))
            return " ".join(parts)

        def table_rows(table: Table, depth: int) -> Iterator[str]:
            nonlocal table_cell_count
            if depth > _MAX_DOCX_TABLE_DEPTH:
                raise DocxSafetyError(
                    f"table nesting exceeds {_MAX_DOCX_TABLE_DEPTH} levels"
                )
            for row in table.rows:
                cells: list[str] = []
                # row.cells expands w:gridSpan before callers can apply a bound.
                # Physical w:tc elements keep malformed documents proportional to XML size.
                for tc in row._tr.tc_lst:
                    table_cell_count += 1
                    if table_cell_count > _MAX_DOCX_TABLE_CELLS:
                        raise DocxSafetyError(
                            f"document contains more than {_MAX_DOCX_TABLE_CELLS} table cells"
                        )
                    cells.append(cell_text(_Cell(tc, table), depth))
                if any(cells):
                    yield "\t".join(cells)

        for block in doc.iter_inner_content():
            if isinstance(block, Paragraph):
                text = block.text.strip()
                if text and not collector.add(text, separator="\n\n"):
                    break
                continue
            if not isinstance(block, Table):
                continue
            first_row = True
            for row_text in table_rows(block, 1):
                separator = "\n\n" if first_row else "\n"
                first_row = False
                if not collector.add(row_text, separator=separator):
                    return collector.render()
        return collector.render()
    except DocxSafetyError as e:
        return f"[error: unsafe DOCX: {e!s}]"
    except Exception as e:
        logger.exception("Failed to extract DOCX {}", path)
        return f"[error: failed to extract DOCX: {e!s}]"


def _extract_xlsx(path: Path) -> str:
    """Extract text from XLSX using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "[error: openpyxl not installed]"
    try:
        if error := _office_archive_error(path):
            return error
        wb = load_workbook(path, read_only=True, data_only=True)
        try:
            collector = _TextCollector(_MAX_TEXT_LENGTH)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                wrote_header = False
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        if not wrote_header:
                            if not collector.add(
                                f"--- Sheet: {sheet_name} ---",
                                separator="\n\n",
                            ):
                                return collector.render()
                            wrote_header = True
                        if not collector.add(row_text, separator="\n"):
                            return collector.render()
            return collector.render()
        finally:
            wb.close()
    except Exception as e:
        logger.exception("Failed to extract XLSX {}", path)
        return f"[error: failed to extract XLSX: {e!s}]"


def _extract_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation as PptxPresentation
    except ImportError:
        return "[error: python-pptx not installed]"
    try:
        if error := _office_archive_error(path):
            return error
        prs = PptxPresentation(path)
        collector = _TextCollector(_MAX_TEXT_LENGTH)
        for i, slide in enumerate(prs.slides, 1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                _collect_pptx_shape_text(shape, slide_text)
            if slide_text:
                if not collector.add(
                    f"--- Slide {i} ---\n" + "\n".join(slide_text),
                    separator="\n\n",
                ):
                    break
        return collector.render()
    except Exception as e:
        logger.exception("Failed to extract PPTX {}", path)
        return f"[error: failed to extract PPTX: {e!s}]"


def _collect_pptx_shape_text(shape, out: list[str]) -> None:
    """Collect text from a PPTX shape, recursing into groups and tables.

    Groups have ``has_text_frame=False`` and must be walked via ``.shapes``;
    tables are GraphicFrame objects whose cell text lives under ``.table``.
    """
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub in sub_shapes:
            _collect_pptx_shape_text(sub, out)
        return

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(cell for cell in cells if cell)
            if line:
                out.append(line)
        return

    text = getattr(shape, "text", "")
    if text:
        out.append(text)


def _office_archive_error(path: Path) -> str | None:
    """Reject oversized or encrypted OOXML containers before parsing XML."""
    try:
        with ZipFile(path) as archive:
            members = archive.infolist()
    except (BadZipFile, OSError) as e:
        return f"[error: invalid Office document: {e!s}]"
    if len(members) > _MAX_OFFICE_ARCHIVE_MEMBERS:
        return f"[error: Office document contains too many files ({len(members)})]"
    total_size = 0
    for member in members:
        if member.flag_bits & 0x1:
            return "[error: encrypted Office documents are not supported]"
        if member.file_size > _MAX_OFFICE_MEMBER_SIZE:
            return "[error: Office document contains an oversized internal file]"
        total_size += member.file_size
        if total_size > _MAX_OFFICE_UNCOMPRESSED_SIZE:
            limit_mb = _MAX_OFFICE_UNCOMPRESSED_SIZE / (1024 * 1024)
            return f"[error: Office document expands beyond the {limit_mb:g} MB safety limit]"
    return None


def _extract_text_file(path: Path) -> str:
    """Extract text from a plain text file."""
    try:
        # Try UTF-8 first, then latin-1 fallback
        try:
            content = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            content = path.read_text(encoding="latin-1")
        return _truncate(content, _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.exception("Failed to read text file {}", path)
        return f"[error: failed to read file: {e!s}]"


def _truncate(text: str, max_length: int) -> str:
    """Truncate text with a suffix indicating truncation."""
    if len(text) <= max_length:
        return text
    return text[:max_length] + f"... (truncated, {len(text)} chars total)"


def _is_text_extension(ext: str) -> bool:
    """Check if extension is a text format."""
    return ext in {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".log",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
        ".cfg",
    }


# ---------------------------------------------------------------------------
# High-level helper: split media into images + extracted document text
# ---------------------------------------------------------------------------


def is_image_file(path: str) -> bool:
    """Check whether *path* looks like an image file.

    Uses magic-byte detection (reads first 16 bytes) with a ``mimetypes``
    extension-based fallback.
    """
    p = Path(path)
    mime: str | None = None
    if p.is_file():
        try:
            with p.open("rb") as f:
                mime = detect_image_mime(f.read(16))
        except OSError:
            mime = None
    if not mime:
        mime = mimetypes.guess_type(path)[0]
    return bool(mime and mime.startswith("image/"))


def reference_non_image_attachments(
    content: str, media: list[str],
) -> tuple[str, list[str]]:
    """Separate images from non-image attachments without reading file content.

    Image paths are preserved for downstream vision-block construction.
    Non-image paths are appended as ``[Attachment: path]`` references.
    """
    image_paths: list[str] = []
    attachment_refs: list[str] = []
    for path in media:
        if is_image_file(path):
            image_paths.append(path)
        else:
            attachment_refs.append(f"[Attachment: {path}]")
    if attachment_refs:
        suffix = "\n".join(attachment_refs)
        content = f"{content}\n\n{suffix}" if content else suffix
    return content, image_paths


def extract_documents(
    text: str,
    media_paths: list[str],
    *,
    max_file_size: int = _MAX_EXTRACT_FILE_SIZE,
) -> tuple[str, list[str]]:
    """Separate images from documents in *media_paths*.

    Documents (PDF, DOCX, XLSX, PPTX, plain-text, …) have their text
    extracted and appended to *text*.  Only image paths are kept in the
    returned list so that downstream layers only need to handle vision
    blocks.

    Files larger than *max_file_size* bytes are skipped with a warning
    to avoid unbounded memory / CPU usage.
    """
    image_paths: list[str] = []
    doc_texts: list[str] = []

    for path_str in media_paths:
        p = Path(path_str)
        if not p.is_file():
            continue

        try:
            size = p.stat().st_size
        except OSError:
            continue
        if size > max_file_size:
            logger.warning(
                "Skipping oversized file for extraction: {} ({:.1f} MB > {} MB limit)",
                p.name, size / (1024 * 1024), max_file_size // (1024 * 1024),
            )
            continue

        if is_image_file(path_str):
            image_paths.append(path_str)
        else:
            extracted = extract_text(p)
            if extracted and not extracted.startswith("[error:"):
                doc_texts.append(f"[File: {p.name}]\n{extracted}")

    if doc_texts:
        text = text + "\n\n" + "\n\n".join(doc_texts)

    return text, image_paths
